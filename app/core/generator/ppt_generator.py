from pptx import Presentation
from pptx.util import Inches
import os
from app.config import settings
from app.core.generator.llm_client import get_llm_client

class PPTGenerator:
    def __init__(self):
        """初始化PPT生成器"""
        self.llm_client = get_llm_client()
    
    def generate_ppt(self, topic: str, content: list, output_filename: str = None, template_id: str = None):
        """
        生成PPT文件
        
        Args:
            topic: PPT主题
            content: PPT内容列表，每个元素包含标题和内容
            output_filename: 输出文件名
            template_id: 模板ID
            
        Returns:
            str: 生成的PPT文件路径
        """
        # 创建PPT对象
        prs = Presentation()
        
        # 添加标题页
        self._add_title_slide(prs, topic)
        
        # 添加内容页
        for item in content:
            if isinstance(item, dict) and 'title' in item and 'content' in item:
                self._add_content_slide(prs, item['title'], item['content'])
        
        # 应用模板（暂时不支持模板ID，使用外部模板文件）
        if template_id:
            print(f"模板ID: {template_id}（暂时不支持）")
        
        # 生成文件名
        if not output_filename:
            output_filename = f"{topic.replace(' ', '_')}_ppt.pptx"
        
        # 确保输出目录存在
        os.makedirs(settings.EXPORT_DIR, exist_ok=True)
        
        # 保存PPT文件
        output_path = os.path.join(settings.EXPORT_DIR, output_filename)
        prs.save(output_path)
        
        return output_path
    
    def _add_title_slide(self, prs, title):
        """添加标题页"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
    
    def _add_content_slide(self, prs, title, content):
        """添加内容页，支持自动分页"""
        # 内容分页阈值（每页最多显示的内容项数）
        PAGE_SIZE = 6
        
        if isinstance(content, list):
            # 计算需要的页数
            total_items = len(content)
            total_pages = (total_items + PAGE_SIZE - 1) // PAGE_SIZE
            
            for page in range(total_pages):
                # 计算当前页的内容范围
                start_idx = page * PAGE_SIZE
                end_idx = min((page + 1) * PAGE_SIZE, total_items)
                page_content = content[start_idx:end_idx]
                
                # 创建幻灯片
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                # 设置标题（添加页码标识）
                title_placeholder = slide.shapes.title
                if total_pages > 1:
                    slide_title = f"{title}（{page + 1}/{total_pages}）"
                else:
                    slide_title = title
                title_placeholder.text = slide_title
                
                # 设置内容
                content_placeholder = slide.placeholders[1]
                content_text = '\n'.join(page_content)
                content_placeholder.text = content_text
        else:
            # 非列表内容直接添加
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            # 设置标题
            title_placeholder = slide.shapes.title
            title_placeholder.text = title
            
            # 设置内容
            content_placeholder = slide.placeholders[1]
            content_text = str(content)
            content_placeholder.text = content_text
    
    def generate_ppt_from_outline(self, outline: dict, output_filename: str = None, template_id: str = None):
        """
        从大纲生成PPT
        
        Args:
            outline: PPT大纲，包含主题和章节
            output_filename: 输出文件名
            template_id: 模板ID
            
        Returns:
            str: 生成的PPT文件路径
        """
        topic = outline.get('topic', 'Untitled')
        sections = outline.get('sections', [])
        
        content = []
        for section in sections:
            section_title = section.get('title', 'Section')
            section_content = section.get('content', [])
            content.append({
                'title': section_title,
                'content': section_content
            })
        
        return self.generate_ppt(topic, content, output_filename, template_id)
    
    def generate_ppt_with_ai(self, topic: str, output_filename: str = None):
        """
        使用AI生成完整PPT

        Args:
            topic: PPT主题
            output_filename: 输出文件名

        Returns:
            str: 生成的PPT文件路径
        """
        # 使用AI生成PPT内容
        outline = self._generate_ppt_outline_with_ai(topic, None)

        # 提取内容
        content = []
        for section in outline.get('sections', []):
            content.append({
                'title': section.get('title', '无标题'),
                'content': section.get('content', [])
            })

        # 随机获取模板
        template = self.get_random_template()

        if template:
            # 使用模板生成PPT
            print(f"使用模板: {template['name']}")
            return self.generate_ppt_from_template(template['path'], topic, content, output_filename)
        else:
            # 如果没有可用模板，使用默认方式生成
            print("没有可用模板，使用默认方式生成")
            return self.generate_ppt(topic, content, output_filename)

    def generate_ppt_from_lesson_plan(self, topic: str, lesson_plan_data: dict, output_filename: str = None):
        """
        从教案数据生成配套PPT

        Args:
            topic: PPT主题
            lesson_plan_data: 教案数据，包含教学目标、教学过程等内容
            output_filename: 输出文件名

        Returns:
            str: 生成的PPT文件路径
        """
        content = []

        # 从教学目标生成内容
        teaching_objectives = lesson_plan_data.get('teaching_objectives', [])
        if teaching_objectives:
            content.append({
                'title': '教学目标',
                'content': teaching_objectives if isinstance(teaching_objectives, list) else [teaching_objectives]
            })

        # 从教学过程生成内容，增强为知识点和详细解释
        teaching_process = lesson_plan_data.get('teaching_process', [])
        if teaching_process:
            for i, process in enumerate(teaching_process, 1):
                if isinstance(process, dict):
                    title = process.get('title', f'教学环节{i}')
                    process_content = process.get('content', [])
                    duration = process.get('duration', '')
                    
                    # 增强内容，确保每个知识点都有详细解释
                    enhanced_content = []
                    if isinstance(process_content, list):
                        for item in process_content:
                            if isinstance(item, str):
                                # 尝试将内容解析为知识点和解释
                                if ':' in item:
                                    enhanced_content.append(item)
                                else:
                                    # 如果没有解释，添加默认解释格式
                                    enhanced_content.append(f"{item}：详细解释")
                    else:
                        enhanced_content = [f"{process_content}：详细解释"]
                    
                    if duration:
                        enhanced_content = [f"{t}（{duration}）" for t in enhanced_content]
                    
                    content.append({
                        'title': title,
                        'content': enhanced_content
                    })
                elif isinstance(process, str):
                    content.append({
                        'title': f'教学环节{i}',
                        'content': [f"{process}：详细解释"]
                    })

        # 从课堂活动生成内容
        classroom_activities = lesson_plan_data.get('classroom_activities', [])
        if classroom_activities:
            activities_content = []
            for activity in classroom_activities:
                if isinstance(activity, dict):
                    activity_title = activity.get('title', '活动')
                    activity_desc = activity.get('description', '')
                    duration = activity.get('duration', '')
                    if activity_desc:
                        activities_content.append(f"{activity_title}：{activity_desc}" + (f"（{duration}）" if duration else ""))
                    else:
                        activities_content.append(f"{activity_title}：详细活动说明" + (f"（{duration}）" if duration else ""))
                elif isinstance(activity, str):
                    activities_content.append(f"{activity}：详细活动说明")
            if activities_content:
                content.append({
                    'title': '课堂活动',
                    'content': activities_content
                })

        # 从课后作业生成内容
        homework = lesson_plan_data.get('homework', [])
        if homework:
            homework_content = []
            if isinstance(homework, list):
                for item in homework:
                    if isinstance(item, str):
                        homework_content.append(f"{item}：完成要求")
                    else:
                        homework_content.append(f"{str(item)}：完成要求")
            else:
                homework_content = [f"{homework}：完成要求"]
            content.append({
                'title': '课后作业',
                'content': homework_content
            })

        # 增强内容，确保包含完整的知识点和解释
        enhanced_content = []
        for slide in content:
            slide_title = slide.get('title', '无标题')
            slide_content = slide.get('content', [])
            
            # 使用AI增强内容，确保每个知识点都有详细解释
            try:
                enhanced_slide_content = self._enhance_slide_content_with_ai(topic, slide_title, slide_content)
                enhanced_content.append({
                    'title': slide_title,
                    'content': enhanced_slide_content
                })
            except Exception:
                # 如果增强失败，使用原始内容
                enhanced_content.append(slide)

        # 如果没有内容，生成默认内容
        if not enhanced_content:
            enhanced_content = [
                {'title': '教学目标', 'content': [f'了解{topic}的基本概念：掌握{topic}的定义和重要性', f'理解{topic}的核心原理：深入学习{topic}的基本原理', f'应用{topic}的知识：能够在实际场景中应用{topic}的知识']},
                {'title': '核心知识点', 'content': [f'{topic}的定义：详细解释{topic}的概念和内涵', f'{topic}的特点：分析{topic}的主要特征和优势', f'{topic}的应用：介绍{topic}在各个领域的应用']},
                {'title': '教学内容', 'content': [f'知识点1：详细解释和案例分析', f'知识点2：详细解释和案例分析', f'知识点3：详细解释和案例分析']},
                {'title': '课堂活动', 'content': ['互动讨论：围绕核心知识点展开讨论', '小组合作：完成相关任务和练习', '案例分析：分析实际应用案例']},
                {'title': '课后作业', 'content': ['复习本节内容：巩固所学知识点', '完成练习题：检验学习效果', '拓展阅读：深入了解相关知识']}
            ]
        else:
            # 确保内容足够丰富，添加核心知识点章节
            has_core_knowledge = any('核心知识' in slide['title'] or '知识点' in slide['title'] for slide in enhanced_content)
            if not has_core_knowledge:
                # 添加核心知识点章节
                core_knowledge_slide = {
                    'title': '核心知识点',
                    'content': [f'{topic}的定义：详细解释{topic}的概念和内涵', f'{topic}的特点：分析{topic}的主要特征和优势', f'{topic}的应用：介绍{topic}在各个领域的应用']
                }
                # 使用AI增强核心知识点内容
                try:
                    core_knowledge_slide['content'] = self._enhance_slide_content_with_ai(topic, '核心知识点', core_knowledge_slide['content'])
                except Exception:
                    pass
                # 插入到内容列表的合适位置
                if enhanced_content:
                    enhanced_content.insert(1, core_knowledge_slide)
                else:
                    enhanced_content.append(core_knowledge_slide)

        # 随机获取模板
        template = self.get_random_template()

        if template:
            print(f"使用模板: {template['name']}")
            return self.generate_ppt_from_template(template['path'], topic, enhanced_content, output_filename)
        else:
            print("没有可用模板，使用默认方式生成")
            return self.generate_ppt(topic, enhanced_content, output_filename)
    
    def enhance_ppt_content_with_ai(self, topic: str, content: list, output_filename: str = None, template_id: str = None):
        """
        使用AI增强PPT内容
        
        Args:
            topic: PPT主题
            content: PPT内容列表
            output_filename: 输出文件名
            template_id: 模板ID
            
        Returns:
            str: 生成的PPT文件路径
        """
        # 使用AI增强每个幻灯片的内容
        enhanced_content = []
        for slide in content:
            slide_title = slide.get('title', '无标题')
            slide_content = slide.get('content', [])
            
            # 使用AI增强内容
            enhanced_slide_content = self._enhance_slide_content_with_ai(topic, slide_title, slide_content)
            
            enhanced_content.append({
                'title': slide_title,
                'content': enhanced_slide_content
            })
        
        # 生成PPT
        return self.generate_ppt(topic, enhanced_content, output_filename, template_id)
    
    def generate_ppt_from_template(self, template_path: str, topic: str, content: list = None, output_filename: str = None):
        """
        从外部模板文件生成PPT
        严格沿用模板的版式、字体、配色、排版样式和母版格式
        只保留模板框架和美化，替换所有文字内容
        
        Args:
            template_path: 模板文件路径
            topic: PPT主题
            content: PPT内容列表（可选）
            output_filename: 输出文件名
            
        Returns:
            str: 生成的PPT文件路径
        """
        # 从模板创建PPT
        prs = Presentation(template_path)
        
        # 保留第一页作为标题页，删除其他所有幻灯片
        # 注意：需要从后往前删除，避免索引问题
        for i in range(len(prs.slides) - 1, 0, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[1:]
        
        # 更新标题页
        if prs.slides:
            title_slide = prs.slides[0]
            title_shape = title_slide.shapes.title
            if title_shape:
                title_shape.text = topic
        
        # 如果提供了内容，添加新的幻灯片
        if content:
            # 直接添加新的内容幻灯片，使用模板的布局
            for item in content:
                if isinstance(item, dict) and 'title' in item and 'content' in item:
                    # 内容分页阈值（每页最多显示的内容项数）
                    PAGE_SIZE = 6
                    
                    title = item['title']
                    slide_content = item['content']
                    
                    if isinstance(slide_content, list):
                        # 计算需要的页数
                        total_items = len(slide_content)
                        total_pages = (total_items + PAGE_SIZE - 1) // PAGE_SIZE
                        
                        for page in range(total_pages):
                            # 计算当前页的内容范围
                            start_idx = page * PAGE_SIZE
                            end_idx = min((page + 1) * PAGE_SIZE, total_items)
                            page_content = slide_content[start_idx:end_idx]
                            
                            # 确保使用模板的布局
                            if len(prs.slide_layouts) > 1:
                                slide_layout = prs.slide_layouts[1]  # 使用模板的内容布局
                                slide = prs.slides.add_slide(slide_layout)
                                
                                # 设置标题（添加页码标识）
                                title_placeholder = slide.shapes.title
                                if title_placeholder:
                                    if total_pages > 1:
                                        slide_title = f"{title}（{page + 1}/{total_pages}）"
                                    else:
                                        slide_title = title
                                    title_placeholder.text = slide_title
                                
                                # 设置内容
                                content_placeholder = None
                                # 尝试多种方式找到内容占位符
                                
                                # 方法1：尝试找到类型为BODY或OBJECT的占位符（内容占位符）
                                for placeholder in slide.placeholders:
                                    placeholder_type = placeholder.placeholder_format.type
                                    # 1=TITLE, 2=BODY, 7=OBJECT, 3=CENTER_TITLE, 4=SUBTITLE
                                    if placeholder_type in (2, 7):  # BODY或OBJECT类型
                                        content_placeholder = placeholder
                                        break
                                
                                if not content_placeholder:
                                    # 方法2：尝试找到不是标题的占位符
                                    for placeholder in slide.placeholders:
                                        placeholder_type = placeholder.placeholder_format.type
                                        if placeholder_type not in (1, 3, 4):  # 不是标题类型
                                            content_placeholder = placeholder
                                            break
                                
                                if not content_placeholder:
                                    # 方法3：尝试找到第二个占位符
                                    if len(slide.placeholders) > 1:
                                        content_placeholder = slide.placeholders[1]
                                
                                if content_placeholder:
                                    content_text = '\n'.join(page_content)
                                    content_placeholder.text = content_text
                                else:
                                    print(f"警告：在模板中找不到内容占位符，幻灯片 '{slide_title}' 可能只有标题")
                    else:
                        # 非列表内容直接添加
                        if len(prs.slide_layouts) > 1:
                            slide_layout = prs.slide_layouts[1]  # 使用模板的内容布局
                            slide = prs.slides.add_slide(slide_layout)
                            
                            # 设置标题
                            title_placeholder = slide.shapes.title
                            if title_placeholder:
                                title_placeholder.text = title
                            
                            # 设置内容
                            content_placeholder = None
                            # 尝试多种方式找到内容占位符
                            
                            # 方法1：尝试找到类型为BODY或OBJECT的占位符（内容占位符）
                            for placeholder in slide.placeholders:
                                placeholder_type = placeholder.placeholder_format.type
                                # 1=TITLE, 2=BODY, 7=OBJECT, 3=CENTER_TITLE, 4=SUBTITLE
                                if placeholder_type in (2, 7):  # BODY或OBJECT类型
                                    content_placeholder = placeholder
                                    break
                            
                            if not content_placeholder:
                                # 方法2：尝试找到不是标题的占位符
                                for placeholder in slide.placeholders:
                                    placeholder_type = placeholder.placeholder_format.type
                                    if placeholder_type not in (1, 3, 4):  # 不是标题类型
                                        content_placeholder = placeholder
                                        break
                            
                            if not content_placeholder:
                                # 方法3：尝试找到第二个占位符
                                if len(slide.placeholders) > 1:
                                    content_placeholder = slide.placeholders[1]
                            
                            if content_placeholder:
                                content_text = str(slide_content)
                                content_placeholder.text = content_text
                            else:
                                print(f"警告：在模板中找不到内容占位符，幻灯片 '{title}' 可能只有标题")
        
        # 生成文件名
        if not output_filename:
            output_filename = f"{topic.replace(' ', '_')}_ppt.pptx"
        
        # 确保输出目录存在
        os.makedirs(settings.EXPORT_DIR, exist_ok=True)
        
        # 保存PPT文件
        output_path = os.path.join(settings.EXPORT_DIR, output_filename)
        prs.save(output_path)
        
        return output_path
    
    def generate_ppt_with_ai_from_template(self, template_path: str, topic: str, content_requirement: str = None, output_filename: str = None):
        """
        使用AI生成内容并从模板文件生成PPT
        严格沿用模板的版式、字体、配色、排版样式和母版格式
        
        Args:
            template_path: 模板文件路径
            topic: PPT主题
            content_requirement: 用户的内容需求（可选）
            output_filename: 输出文件名
            
        Returns:
            str: 生成的PPT文件路径
        """
        # 使用AI生成内容
        outline = self._generate_ppt_outline_with_ai(topic, content_requirement)
        
        # 提取内容
        content = []
        for section in outline.get('sections', []):
            content.append({
                'title': section.get('title', '无标题'),
                'content': section.get('content', [])
            })
        
        # 使用模板生成PPT
        return self.generate_ppt_from_template(template_path, topic, content, output_filename)
    
    def generate_enhanced_ppt_with_ai_from_template(self, template_path: str, topic: str, content_requirement: str = None, output_filename: str = None):
        """
        使用AI生成并增强内容，从模板文件生成PPT
        严格沿用模板的版式、字体、配色、排版样式和母版格式
        只生成一个包含AI增强内容的PPT文件
        
        Args:
            template_path: 模板文件路径
            topic: PPT主题
            content_requirement: 用户的内容需求（可选）
            output_filename: 输出文件名
            
        Returns:
            str: 生成的PPT文件路径
        """
        # 使用AI生成内容
        outline = self._generate_ppt_outline_with_ai(topic, content_requirement)
        
        # 提取内容
        content = []
        for section in outline.get('sections', []):
            content.append({
                'title': section.get('title', '无标题'),
                'content': section.get('content', [])
            })
        
        # 使用AI增强内容
        enhanced_content = []
        for slide in content:
            slide_title = slide.get('title', '无标题')
            slide_content = slide.get('content', [])
            
            # 使用AI增强内容
            enhanced_slide_content = self._enhance_slide_content_with_ai(topic, slide_title, slide_content)
            
            enhanced_content.append({
                'title': slide_title,
                'content': enhanced_slide_content
            })
        
        # 使用模板生成PPT
        return self.generate_ppt_from_template(template_path, topic, enhanced_content, output_filename)
    
    def _generate_ppt_outline_with_ai(self, topic: str, content_requirement: str = None):
        """
        使用AI生成PPT大纲
        
        Args:
            topic: PPT主题
            content_requirement: 用户的内容需求（可选）
            
        Returns:
            dict: PPT大纲
        """
        prompt = f"请为'{topic}'主题生成一个完整的PPT大纲，包含以下结构：\n"
        prompt += "{\n"
        prompt += "  \"topic\": \"主题名称\",\n"
        prompt += "  \"sections\": [\n"
        prompt += "    {\n"
        prompt += "      \"title\": \"章节标题\",\n"
        prompt += "      \"content\": [\"知识点1：知识点解释\", \"知识点2：知识点解释\", \"知识点3：知识点解释\"]\n"
        prompt += "    }\n"
        prompt += "  ]\n"
        prompt += "}\n"
        
        if content_requirement:
            prompt += f"\n用户内容需求：{content_requirement}\n"
        
        prompt += "请确保大纲结构完整，内容合理，至少包含4个章节，每个章节至少包含3个知识点，每个知识点都要有详细的解释。PPT内容应包含完整的教学内容，包括核心知识点、详细解释、教学案例等。"
        
        try:
            response = self.llm_client.generate(prompt)
            # 尝试解析响应为JSON
            import json
            try:
                # 去除可能的markdown代码块标记
                if response.startswith("```json"):
                    response = response[7:]
                if response.endswith("```"):
                    response = response[:-3]
                outline = json.loads(response.strip())
                return outline
            except json.JSONDecodeError:
                # 如果解析失败，返回默认大纲
                print("解析AI生成的大纲失败，使用默认大纲")
                return {
                    'topic': topic,
                    'sections': [
                        {
                            'title': f'{topic}概述',
                            'content': [f'{topic}的定义', f'{topic}的重要性', f'{topic}的应用场景']
                        },
                        {
                            'title': f'{topic}的核心内容',
                            'content': ['核心概念', '关键技术', '实现方法']
                        },
                        {
                            'title': f'{topic}的未来发展',
                            'content': ['发展趋势', '挑战与机遇', '总结与展望']
                        }
                    ]
                }
        except Exception as e:
            print(f"生成PPT大纲失败: {str(e)}")
            # 失败时返回默认大纲
            return {
                'topic': topic,
                'sections': [
                    {
                        'title': f'{topic}概述',
                        'content': [f'{topic}的定义', f'{topic}的重要性', f'{topic}的应用场景']
                    },
                    {
                        'title': f'{topic}的核心内容',
                        'content': ['核心概念', '关键技术', '实现方法']
                    },
                    {
                        'title': f'{topic}的未来发展',
                        'content': ['发展趋势', '挑战与机遇', '总结与展望']
                    }
                ]
            }
    
    def get_available_templates(self):
        """
        获取所有可用模板

        Returns:
            list: 模板列表，每个模板包含路径和名称
        """
        templates_dir = os.path.join(os.path.dirname(__file__), 'templates', 'external')
        templates = []

        if os.path.exists(templates_dir):
            for filename in os.listdir(templates_dir):
                if filename.endswith('.pptx'):
                    template_path = os.path.join(templates_dir, filename)
                    templates.append({
                        'id': filename.replace('.pptx', ''),
                        'name': filename.replace('.pptx', '').replace('_', ' '),
                        'path': template_path
                    })

        return templates

    def get_random_template(self):
        """
        随机获取一个可用模板

        Returns:
            dict: 随机选择的模板信息，包含路径和名称
        """
        templates = self.get_available_templates()
        if not templates:
            return None
        import random
        return random.choice(templates)
    

    
    def _enhance_slide_content_with_ai(self, topic: str, slide_title: str, slide_content: list):
        """
        使用AI增强幻灯片内容
        
        Args:
            topic: PPT主题
            slide_title: 幻灯片标题
            slide_content: 幻灯片内容
            
        Returns:
            list: 增强后的内容
        """
        prompt = f"请增强'{topic}'主题下'{slide_title}'幻灯片的内容。\n"
        prompt += f"当前内容：{', '.join(slide_content)}\n"
        prompt += "请生成更详细、更有深度的内容，保持要点式结构，每个要点简洁明了。\n"
        prompt += "请以列表形式输出，每个要点一行。"
        
        try:
            response = self.llm_client.generate(prompt)
            # 解析响应，提取增强后的内容
            enhanced_content = [line.strip() for line in response.split('\n') if line.strip()]
            # 确保内容不为空
            return enhanced_content if enhanced_content else slide_content
        except Exception as e:
            print(f"增强幻灯片内容失败: {str(e)}")
            # 失败时返回原始内容
            return slide_content

# 测试代码
if __name__ == "__main__":
    generator = PPTGenerator()
    
    # 测试基本生成
    test_topic = "人工智能简介"
    test_content = [
        {
            'title': '什么是人工智能',
            'content': ['人工智能是模拟人类智能的技术', '包括机器学习、深度学习等分支', '应用于多个领域']
        },
        {
            'title': '人工智能的历史',
            'content': ['1956年诞生', '经历两次寒冬', '近年来快速发展']
        },
        {
            'title': '人工智能的应用',
            'content': ['自动驾驶', '语音识别', '图像识别', '自然语言处理']
        }
    ]
    
    # 测试使用模板生成
    output_path = generator.generate_ppt(test_topic, test_content, template_id="education")
    print(f"使用教育模板生成PPT成功: {output_path}")
    
    # 测试从大纲生成
    test_outline = {
        'topic': '机器学习基础',
        'sections': [
            {
                'title': '机器学习概述',
                'content': ['定义', '类型', '应用场景']
            },
            {
                'title': '监督学习',
                'content': ['分类', '回归', '常见算法']
            },
            {
                'title': '无监督学习',
                'content': ['聚类', '降维', '异常检测']
            }
        ]
    }
    
    output_path2 = generator.generate_ppt_from_outline(test_outline, template_id="technology")
    print(f"使用科技模板从大纲生成PPT成功: {output_path2}")
    
    # 测试获取可用模板
    templates = generator.get_available_templates()
    print(f"可用模板: {templates}")
    
    # 测试模板预览
    preview = generator.get_template_preview("business")
    print(f"商务模板预览: {preview}")